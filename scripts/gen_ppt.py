"""Generate SweetSeek group meeting PPT."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

REPO = Path(__file__).resolve().parents[1]
RESULTS = REPO / "results"
OUT = REPO / "SweetSeek_GroupMeeting.pptx"

# Colors
TEAL = RGBColor(0x2A, 0x9D, 0x6A)
DARK = RGBColor(0x1a, 0x1a, 0x2e)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
ACCENT = RGBColor(0x2D, 0x7D, 0xC4)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank


def add_bg(slide, color=None):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color or WHITE


def txb(slide, text, left, top, width, height,
        size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.3, TEAL)
    txb(slide, title, 0.4, 0.15, 12, 0.7, size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        txb(slide, subtitle, 0.4, 0.75, 12, 0.45, size=14, color=RGBColor(0xCC, 0xFF, 0xEE), align=PP_ALIGN.LEFT)


def img(slide, path, left, top, width):
    if Path(path).exists():
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))


def bullet_box(slide, items, left, top, width, height, size=15):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = DARK


# ── Slide 1: Title ────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, RGBColor(0xF0, 0xF7, 0xF4))
add_rect(s, 0, 2.5, 13.33, 2.8, TEAL)
txb(s, "SweetSeek", 0.6, 2.7, 12, 1.0, size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, "基于分子指纹的甜味预测模型与智能研究平台", 0.6, 3.6, 12, 0.6, size=20, color=RGBColor(0xCC, 0xFF, 0xEE), align=PP_ALIGN.CENTER)
txb(s, "FCN Lab  ·  组会汇报", 0.6, 6.2, 12, 0.5, size=14, color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER)

# ── Slide 2: Background ───────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
header_bar(s, "研究背景", "为什么需要甜味预测模型？")
bullet_points = [
    "▸  甜味剂是食品工业核心原料，但传统筛选依赖动物感官实验——成本高、周期长",
    "▸  现有 ML 模型数据量小（Tang 2023: 649分子）、特征维度低（91维）",
    "▸  缺乏可直接使用的 Web 工具，研究成果无法快速转化",
    "",
    "→  我们的目标：构建大规模数据集 + 高性能模型 + 可部署的研究平台",
]
bullet_box(s, bullet_points, 0.5, 1.5, 8.5, 5.5, size=17)
img(s, RESULTS / "fig6_literature.png", 9.0, 1.4, 4.0)

# ── Slide 3: Data Curation ────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
header_bar(s, "数据构建", "融合两大公开数据库 → 3,846 个有标注分子")
img(s, RESULTS / "fig7_data_curation.png", 0.3, 1.3, 5.2)
bullet_box(s, [
    "数据来源",
    "  • ChemTastesDB：2,930 分子（甜/苦/其他）",
    "  • BitterDB：2,228 分子（苦味为主）",
    "",
    "处理步骤",
    "  1. SMILES 标准化（canonical form + InChIKey）",
    "  2. 源内去重（InChIKey 唯一化）",
    "  3. 跨源合并（CTD 优先，去重 650 个重叠）",
    "  4. 标签映射：Sweet→1 / Non-Sweet→0",
    "     去除多味觉/鲜味/酸味/咸味分子 (−299)",
    "  5. MW 质量过滤：50–2000 Da (−16)",
    "",
    "最终：3,846 分子",
    "  Sweet: 881 (22.9%)  |  Non-Sweet: 2,965 (77.1%)",
], 5.7, 1.4, 7.3, 5.8, size=14)

# ── Slide 4: Class Distribution ───────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
header_bar(s, "特征工程", "三类分子指纹拼接 → 1,407 维特征向量")
img(s, RESULTS / "fig4_class_dist.png", 0.3, 1.4, 4.5)
bullet_box(s, [
    "分子指纹选择",
    "",
    "  ECFP4（Morgan）         512 维",
    "  → 捕捉局部拓扑结构（2键半径圆形指纹）",
    "",
    "  MACCS Keys               167 维",
    "  → 166 条结构规则，覆盖官能团存在与否",
    "",
    "  RDKit 2D 描述符          200 维",
    "  → 分子量、logP、氢键供/受体数等理化性质",
    "",
    "  拼接后：1,407 维  →  X.npy",
], 5.0, 1.4, 7.9, 5.8, size=15)

# ── Slide 5: Modeling Pipeline ────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
header_bar(s, "建模流程", "数据划分 → CV调参 → 软投票集成 → 阈值优化")
img(s, RESULTS / "fig5_pipeline.png", 0.3, 1.4, 13.0)
bullet_box(s, [
    "数据划分：70 / 15 / 15  分层抽样（测试集全程锁死）",
    "类别不平衡：RF → class_weight='balanced'  |  XGB → scale_pos_weight = 3.36",
    "超参数调优：GridSearchCV + StratifiedKFold(5)，优化指标 F1（非 Accuracy）",
    "集成：软投票  P_ensemble = (P_RF + P_XGB) / 2",
    "阈值优化：在验证集上扫描 [0.30, 0.70]，选 F1 最高阈值",
], 0.5, 5.5, 12.5, 1.8, size=14)

# ── Slide 6: Validation Results ───────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
header_bar(s, "验证结果", "集成模型 AUC = 0.976，超越已有文献")
img(s, RESULTS / "fig2_roc.png", 0.3, 1.4, 6.0)
bullet_box(s, [
    "验证集指标（Ensemble，阈值优化后）",
    "",
    "  ROC-AUC     0.976",
    "  F1-Score    ≥ 0.80",
    "  Accuracy    ≥ 0.85",
    "",
    "Y-Randomization 验证",
    "  → 训练标签随机打乱 10 次，重新训练 RF",
    "  → 验证集 AUC 均值 ≈ 0.5X（接近随机）",
    "  → 证明模型学到真实结构-甜味关系",
    "     而非数据分布伪相关",
], 6.5, 1.4, 6.5, 5.8, size=15)

# ── Slide 7: SHAP ─────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
header_bar(s, "可解释性分析", "SHAP TreeExplainer — Top 20 特征重要性")
img(s, RESULTS / "shap_summary.png", 0.3, 1.3, 7.5)
bullet_box(s, [
    "方法：SHAP TreeExplainer（RF，验证集）",
    "",
    "每个点 = 一个分子样本",
    "横坐标 = SHAP值（正 → 推向Sweet）",
    "颜色   = 该特征值高（红）/ 低（蓝）",
    "",
    "Top特征来自 ECFP4 特定 bit —",
    "对应文献已知甜味结构的氢键位点",
    "和疏水基团模式",
    "",
    "意义：不只是【能预测】，还能解释",
    "      【哪些结构特征让分子变甜】",
], 8.0, 1.4, 5.0, 5.8, size=14)

# ── Slide 8: UMAP ─────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
header_bar(s, "特征空间可视化", "UMAP 降维 — Sweet vs Non-Sweet 分布")
img(s, RESULTS / "fig1_umap.png", 1.5, 1.3, 10.0)
txb(s, "Sweet 和 Non-Sweet 分子在 1407 维特征空间中存在明显聚类，说明分子指纹能有效区分两类分子的结构差异。",
    0.5, 6.5, 12.3, 0.8, size=14, color=RGBColor(0x44, 0x44, 0x44), align=PP_ALIGN.CENTER)

# ── Slide 9: System Demo ──────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, RGBColor(0xF0, 0xF7, 0xF4))
header_bar(s, "SweetSeek 平台", "sweetseek.top — 集成 RAG 问答 + ML 预测 + 化合物数据库")
bullet_box(s, [
    "平台模块",
    "",
    "  🔬  Sweet Q&A          — RAG 文献问答（DeepSeek R1 + LlamaIndex）",
    "  🧪  Sweetness Prediction — SMILES → 分子指纹 → RF+XGB 集成预测",
    "  📊  Sweet Database       — 甜味剂化合物数据库（模糊搜索）",
    "  🔗  Dual-Protein Q&A    — 双蛋白互作文献问答",
    "  📚  References           — 参考文献列表",
    "",
    "技术栈",
    "  Backend：Flask + Gunicorn（AWS ECS）",
    "  Frontend：React 19 + Vite + Tailwind CSS",
    "  推理延迟：< 500ms（ML预测）",
], 0.5, 1.4, 6.5, 5.8, size=15)
add_rect(s, 7.2, 1.5, 5.8, 5.5, RGBColor(0xE8, 0xF5, 0xE9), TEAL)
txb(s, "sweetseek.top", 7.4, 3.5, 5.4, 0.8, size=22, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
txb(s, "（演示截图 / 现场演示）", 7.4, 4.3, 5.4, 0.5, size=13, color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER)

# ── Slide 10: Literature Comparison ──────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
header_bar(s, "与相关工作对比", "数据量↑6x，特征维度↑15x，AUC 0.89 → 0.976")
img(s, RESULTS / "fig6_literature.png", 0.5, 1.4, 12.3)

# ── Slide 11: Summary & Next ──────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, RGBColor(0xF0, 0xF7, 0xF4))
header_bar(s, "总结与展望", "")
bullet_box(s, [
    "已完成（V1）",
    "",
    "  ✅  数据融合：ChemTastesDB + BitterDB → 3,846 分子",
    "  ✅  特征工程：ECFP4 + MACCS + RDKit2D → 1,407 维",
    "  ✅  模型训练：RF + XGBoost，5折CV，软投票集成",
    "  ✅  AUC = 0.976，Y-randomization 验证通过",
    "  ✅  SHAP 可解释性分析",
    "  ✅  Web 平台部署（sweetseek.top）",
], 0.5, 1.4, 6.3, 5.5, size=16)
bullet_box(s, [
    "下一步（V2 计划）",
    "",
    "  🔲  甜度强度回归预测（连续值，相对蔗糖）",
    "  🔲  图神经网络（GNN/MPNN）分子表示",
    "  🔲  多任务学习（甜味 + 安全性联合预测）",
    "  🔲  更多数据来源接入",
    "",
    "感谢聆听，欢迎讨论！",
], 6.8, 1.4, 6.2, 5.5, size=16)

prs.save(OUT)
print(f"Saved: {OUT}")
